#!/usr/bin/env python3
"""
Convert Slide markdown to editable PPTX with CSS-matched styling.

Design approach: Use PPTX-native features to match the CSS visual style.
- Background: pre-rendered image (gradient + accents)
- Cards: rounded rectangles with shadow
- Images: fill full grid cell (no side whitespace)
- Text: exact font/color/size from CSS em chain

Usage:
  python md_to_pptx.py --input slides.md --output slides.pptx \
      --bg-image backgrounds/slide_bg_light.png --base-dir /path/to/project
"""

import re
import os
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
from PIL import Image

# --- Slide canvas: 1280x720 CSS pixels ---
PX = 9525  # EMU per CSS pixel

def px(n):
    return Emu(int(n * PX))

# CSS padding: 40px 36px 50px 36px
PAD_TOP, PAD_RIGHT, PAD_BOTTOM, PAD_LEFT = 40, 36, 50, 36
CANVAS_W, CANVAS_H = 1280, 720
CL = PAD_LEFT  # content left
CW = CANVAS_W - PAD_LEFT - PAD_RIGHT  # 1208px content width

# --- Font sizes (points) ---
HEADING_LEAD_SIZE = 37       # h1 on title slide
HEADING_LEAD_MIN = 22        # min after auto-scale
HEADING_LEAD_CHARS = 40      # chars before scaling kicks in
HEADING_SIZE = 25            # h2 on content slides
HEADING_MIN = 18             # min after auto-scale
HEADING_CHARS = 55           # chars before scaling kicks in
SUBTITLE_SIZE = 16           # lead slide subtitle
SUBTITLE_MIN = 11
SUBTITLE_CHARS = 80
BULLET_FONT_SIZE = 16
TABLE_FONT_SIZE = 16
CAPTION_FONT_SIZE = 12
CHIP_FONT_SIZE = 8
CHIP_LEAD_FONT_SIZE = 12     # chips on title slide

CODE_SIZE_OFFSET = 2         # code font = base - this

# --- Bullet layout (CSS pixels) ---
# BULLET_FONT_SIZE pt * 1.15 PPTX line spacing = 25 CSS px per wrapped line
# space_before(3pt) + space_after(3pt) = 8 CSS px per paragraph
BULLET_LINE_H = 25
BULLET_PARA_PAD = 8
BULLET_SPACE_PT = 3          # space_before and space_after in points

# --- Video support ---
VIDEO_EXTENSIONS = ('.mp4', '.avi', '.mov', '.webm', '.wmv', '.m4v', '.mkv', '.flv')
VIDEO_CONTENT_TYPES = {
    '.mp4': 'video/mp4',
    '.avi': 'video/x-msvideo',
    '.mov': 'video/quicktime',
    '.webm': 'video/webm',
    '.wmv': 'video/x-ms-wmv',
    '.m4v': 'video/x-m4v',
    '.mkv': 'video/x-matroska',
    '.flv': 'video/x-flv',
}
_video_counter = 0

# --- Color palettes ---
# Light palette (default — warm paper and off-white backgrounds)
LIGHT_COLORS = dict(
    INK=RGBColor(0x1D, 0x23, 0x30),
    INK_SOFT=RGBColor(0x4A, 0x55, 0x6F),
    CARD_BG=RGBColor(0xFF, 0xFF, 0xFF),
    LINE_COLOR=RGBColor(0xC8, 0xCA, 0xD0),
    ACCENT_A=RGBColor(0xE4, 0x57, 0x2E),
    ACCENT_B=RGBColor(0x0C, 0x7C, 0x99),
    ACCENT_C=RGBColor(0x58, 0x81, 0x57),
    ACCENT_D=RGBColor(0xF4, 0xA2, 0x59),
    CAPTION_CLR=RGBColor(0x5F, 0x68, 0x81),
    TBL_HDR_BG=RGBColor(0xE8, 0xF4, 0xF7),
    TBL_HDR_TXT=RGBColor(0x22, 0x32, 0x4A),
    TBL_BODY_TXT=RGBColor(0x4A, 0x55, 0x6F),
    CHIP_CLR=RGBColor(0x36, 0x41, 0x55),
    CHIP_BG=RGBColor(0xF2, 0xF2, 0xF6),
    CHIP_BORDER=RGBColor(0xCB, 0xCD, 0xD2),
    STRONG_CLR=RGBColor(0x27, 0x32, 0x48),
    CODE_CLR=RGBColor(0x26, 0x4D, 0x62),
    CARD_ALPHA='86000',
    SHADOW_COLOR='1D2330',
    SHADOW_ALPHA='10000',
)

# Dark palette
DARK_COLORS = dict(
    INK=RGBColor(0xED, 0xEF, 0xF3),          # light text on dark bg
    INK_SOFT=RGBColor(0xD0, 0xD4, 0xDC),      # readable light text
    CARD_BG=RGBColor(0x2C, 0x2C, 0x2E),        # dark card
    LINE_COLOR=RGBColor(0x3A, 0x3A, 0x3C),      # subtle border on dark
    ACCENT_A=RGBColor(0xF0, 0x72, 0x4B),        # brighter coral
    ACCENT_B=RGBColor(0x28, 0xA0, 0xBE),        # brighter teal
    ACCENT_C=RGBColor(0x7A, 0xAA, 0x78),        # brighter sage
    ACCENT_D=RGBColor(0xFA, 0xB8, 0x6C),        # brighter amber
    CAPTION_CLR=RGBColor(0xC8, 0xCC, 0xD4),
    TBL_HDR_BG=RGBColor(0xE8, 0xF4, 0xF7),      # light header bg
    TBL_HDR_TXT=RGBColor(0x1D, 0x23, 0x30),    # black header text
    TBL_BODY_TXT=RGBColor(0x1D, 0x23, 0x30),   # black body text
    CHIP_CLR=RGBColor(0xC8, 0xCE, 0xDA),
    CHIP_BG=RGBColor(0x32, 0x38, 0x4C),
    CHIP_BORDER=RGBColor(0x4A, 0x50, 0x64),
    STRONG_CLR=RGBColor(0xE8, 0xEC, 0xF2),
    CODE_CLR=RGBColor(0x6C, 0xC0, 0xD8),        # bright teal for code
    CARD_ALPHA='80000',
    SHADOW_COLOR='000000',
    SHADOW_ALPHA='20000',
)

# Active palette — set by apply_theme()
INK = INK_SOFT = CARD_BG = LINE_COLOR = None
ACCENT_A = ACCENT_B = ACCENT_C = ACCENT_D = None
CAPTION_CLR = TBL_HDR_BG = TBL_HDR_TXT = TBL_BODY_TXT = None
CHIP_CLR = CHIP_BG = CHIP_BORDER = STRONG_CLR = CODE_CLR = None
CARD_ALPHA = SHADOW_COLOR = SHADOW_ALPHA = None


def apply_theme(theme_name):
    """Set module-level color globals from a theme palette."""
    global INK, INK_SOFT, CARD_BG, LINE_COLOR
    global ACCENT_A, ACCENT_B, ACCENT_C, ACCENT_D
    global CAPTION_CLR, TBL_HDR_BG, TBL_HDR_TXT, TBL_BODY_TXT
    global CHIP_CLR, CHIP_BG, CHIP_BORDER, STRONG_CLR, CODE_CLR
    global CARD_ALPHA, SHADOW_COLOR, SHADOW_ALPHA
    palette = DARK_COLORS if theme_name == 'dark' else LIGHT_COLORS
    for k, v in palette.items():
        globals()[k] = v


# Default to light
apply_theme('light')

FONT_BODY = 'Manrope'
FONT_HEAD = 'Fraunces'
FONT_CODE = 'Consolas'

# Configurable at runtime via CLI args
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
BG_IMAGE_PATH = None  # Set by CLI


# ═══════════════════════════════════════════════════════════════
# Background — rendered from CSS as a single image
# ═══════════════════════════════════════════════════════════════

def set_slide_background(slide, prs):
    """Add background as a full-slide image shape (copies correctly on duplicate)."""
    if BG_IMAGE_PATH is None or not os.path.exists(BG_IMAGE_PATH):
        return

    pic = slide.shapes.add_picture(
        BG_IMAGE_PATH, Emu(0), Emu(0),
        prs.slide_width, prs.slide_height)

    # Send to back: move pic element right after spTree's grpSpPr
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    sp = pic._element
    spTree.remove(sp)
    grpSpPr = spTree.find(qn('p:grpSpPr'))
    grpSpPr.addnext(sp)

    # Lock shape so it can't be selected or moved
    cNvPicPr = sp.find(qn('p:nvPicPr')).find(qn('p:cNvPicPr'))
    picLocks = cNvPicPr.find(qn('a:picLocks'))
    if picLocks is None:
        picLocks = etree.SubElement(cNvPicPr, qn('a:picLocks'))
    picLocks.set('noSelect', '1')
    picLocks.set('noMove', '1')
    picLocks.set('noResize', '1')
    picLocks.set('noRot', '1')
    picLocks.set('noChangeAspect', '1')


# ═══════════════════════════════════════════════════════════════
# Text helpers
# ═══════════════════════════════════════════════════════════════

def _make_txbox(slide, l, t, w, h):
    txBox = slide.shapes.add_textbox(px(l), px(t), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return txBox, tf


def add_text(slide, l, t, w, h, text, font=FONT_BODY, size=13, color=INK_SOFT,
             bold=False, align=PP_ALIGN.LEFT):
    txBox, tf = _make_txbox(slide, l, t, w, h)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_rich(slide, l, t, w, h, text, font=FONT_BODY, size=13, color=INK_SOFT,
             bold=False, align=PP_ALIGN.LEFT):
    """Text with **bold** and `code` inline markup."""
    txBox, tf = _make_txbox(slide, l, t, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    _fill_runs(p, text, font, size, color, bold)
    return txBox


def _fill_runs(p, text, font, size, color, bold):
    parts = re.split(r'(\*\*.*?\*\*|`[^`]+`)', text)
    for part in parts:
        run = p.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = STRONG_CLR
            run.font.bold = True
        elif part.startswith('`') and part.endswith('`'):
            run.text = part[1:-1]
            run.font.name = FONT_CODE
            run.font.size = Pt(max(8, size - CODE_SIZE_OFFSET))
            run.font.color.rgb = CODE_CLR
        else:
            run.text = part
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold


# ═══════════════════════════════════════════════════════════════
# Markdown extractors
# ═══════════════════════════════════════════════════════════════

def extract_heading(t):
    m = re.search(r'^#{1,2}\s+(.+)$', t, re.MULTILINE)
    return m.group(1) if m else None

def extract_bullets(t):
    bullets, in_div = [], 0
    for line in t.split('\n'):
        s = line.strip()
        in_div += s.count('<div') - s.count('</div>')
        in_div = max(0, in_div)
        if in_div == 0 and s.startswith('- '):
            bullets.append(s[2:])
    return bullets

def extract_images(t):
    figs = re.findall(r'<figure class="figure-card">.*?<img src="([^"]+)".*?<figcaption>(.*?)</figcaption>.*?</figure>', t, re.DOTALL)
    return [(src, re.sub(r'<[^>]+>', '', cap).strip()) for src, cap in figs]

def extract_grid_class(t):
    m = re.search(r'<div class="cols-(\d)">', t)
    return int(m.group(1)) if m else 2

def extract_table(t):
    hdrs = re.findall(r'<th>(.*?)</th>', t)
    rows = []
    for tr_match in re.finditer(r'<tr>((?:<td>.*?</td>)+)</tr>', t):
        cells = re.findall(r'<td>(.*?)</td>', tr_match.group(1))
        if cells:
            rows.append(tuple(cells))
    return hdrs, rows

def extract_chips(t):
    return re.findall(r'<span class="chip">(.*?)</span>', t)



# ═══════════════════════════════════════════════════════════════
# Card with shadow (PPTX native shadow matching CSS box-shadow)
# ═══════════════════════════════════════════════════════════════

def add_card(slide, l, t, w, h, fill=CARD_BG, border=LINE_COLOR, shadow=True, radius=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(l), px(t), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    # Set fill transparency for card-bg
    if fill == CARD_BG:
        srgb_el = shp._element.find(qn('p:spPr')).find(qn('a:solidFill')).find(qn('a:srgbClr'))
        alpha_el = etree.SubElement(srgb_el, qn('a:alpha'))
        alpha_el.set('val', CARD_ALPHA)
    shp.line.color.rgb = border
    shp.line.width = Pt(0.75)
    adj = min(0.5, radius / min(w, h)) if min(w, h) > 0 else 0.05
    shp.adjustments[0] = adj

    if shadow:
        sp = shp._element.find(qn('p:spPr'))
        effect_lst = etree.SubElement(sp, qn('a:effectLst'))
        outer = etree.SubElement(effect_lst, qn('a:outerShdw'))
        outer.set('blurRad', str(int(30 * 12700)))
        outer.set('dist', str(int(14 * 12700)))
        outer.set('dir', '5400000')
        outer.set('algn', 'bl')
        srgb = etree.SubElement(outer, qn('a:srgbClr'))
        srgb.set('val', SHADOW_COLOR)
        a = etree.SubElement(srgb, qn('a:alpha'))
        a.set('val', SHADOW_ALPHA)

    return shp


# ═══════════════════════════════════════════════════════════════
# Image placement — fill cell, object-fit: contain
# ═══════════════════════════════════════════════════════════════

def add_image(slide, img_path, l, t, w, h):
    """Place image filling cell width, object-fit: contain, centered."""
    full = os.path.join(BASE_DIR, img_path)
    if not os.path.exists(full):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(l), px(t), px(w), px(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        shp.line.fill.background()
        return
    try:
        img = Image.open(full)
        iw, ih = img.size
        box_r = w / h
        img_r = iw / ih
        if img_r > box_r:
            aw, ah = w, w / img_r
            al, at = l, t + (h - ah) / 2
        else:
            ah, aw = h, h * img_r
            al, at = l + (w - aw) / 2, t
        pic = slide.shapes.add_picture(full, px(al), px(at), px(aw), px(ah))
        # Add rounded corners
        sp = pic._element.find(qn('p:spPr'))
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            sp.remove(prstGeom)
        prstGeom = etree.SubElement(sp, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        corner_ratio = min(10 / min(aw, ah), 0.1) if min(aw, ah) > 0 else 0.02
        gd.set('fmla', f'val {int(corner_ratio * 50000)}')
    except Exception as e:
        print(f"  Warning: {img_path}: {e}")


# ═══════════════════════════════════════════════════════════════
# Video embedding — playable video in PPTX
# ═══════════════════════════════════════════════════════════════

def _is_video(path):
    """Check if path is a video file by extension."""
    return os.path.splitext(path)[1].lower() in VIDEO_EXTENSIONS


def _get_video_dimensions(video_path):
    """Get (width, height) of a video. Returns (1, 1) on failure."""
    try:
        import cv2
        cap = cv2.VideoCapture(video_path)
        w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        cap.release()
        if w > 0 and h > 0:
            return w, h
    except ImportError:
        pass
    try:
        import subprocess
        result = subprocess.run(
            ['ffprobe', '-v', 'quiet', '-select_streams', 'v:0',
             '-show_entries', 'stream=width,height',
             '-of', 'csv=p=0:s=x', video_path],
            capture_output=True, text=True, timeout=10)
        parts = result.stdout.strip().split('x')
        if len(parts) == 2:
            return int(parts[0]), int(parts[1])
    except Exception:
        pass
    return 1, 1


def _extract_poster_frame(video_path):
    """Extract a frame from the middle of a video. Returns temp PNG path or None."""
    import tempfile
    poster_path = tempfile.mktemp(suffix='.png')
    # Try cv2 first
    try:
        import cv2
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            cap.release()
            return None
        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        if total > 0:
            cap.set(cv2.CAP_PROP_POS_FRAMES, total // 2)
        ret, frame = cap.read()
        cap.release()
        if ret:
            cv2.imwrite(poster_path, frame)
            return poster_path
    except ImportError:
        pass
    # Fallback: ffmpeg
    try:
        import subprocess
        result = subprocess.run(
            ['ffprobe', '-v', 'quiet', '-show_entries', 'format=duration',
             '-of', 'default=noprint_wrappers=1:nokey=1', video_path],
            capture_output=True, text=True, timeout=10)
        duration = float(result.stdout.strip()) if result.stdout.strip() else 1.0
        subprocess.run(
            ['ffmpeg', '-y', '-ss', str(duration / 2), '-i', video_path,
             '-vframes', '1', '-q:v', '2', poster_path],
            capture_output=True, timeout=30)
        if os.path.exists(poster_path) and os.path.getsize(poster_path) > 0:
            return poster_path
    except Exception:
        pass
    return None


def add_video(slide, video_path, l, t, w, h):
    """Embed a playable video with auto-generated poster frame in the slide."""
    global _video_counter
    full = os.path.join(BASE_DIR, video_path)

    if not os.path.exists(full):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(l), px(t), px(w), px(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        shp.line.fill.background()
        return

    # Get video dimensions for aspect-correct sizing
    vw, vh = _get_video_dimensions(full)
    box_r = w / h
    vid_r = vw / vh
    if vid_r > box_r:
        aw, ah = w, w / vid_r
        al, at = l, t + (h - ah) / 2
    else:
        ah, aw = h, h * vid_r
        al, at = l + (w - aw) / 2, t

    # Extract poster frame
    poster_path = _extract_poster_frame(full)
    poster_cleanup = poster_path

    if not poster_path:
        # Create a grey placeholder image as poster
        import tempfile
        poster_path = tempfile.mktemp(suffix='.png')
        placeholder = Image.new('RGB', (max(1, vw), max(1, vh)), (0xDD, 0xDD, 0xDD))
        placeholder.save(poster_path)
        poster_cleanup = poster_path

    try:
        # Add poster as picture shape
        pic = slide.shapes.add_picture(poster_path, px(al), px(at), px(aw), px(ah))

        # Add rounded corners (same as add_image)
        sp = pic._element.find(qn('p:spPr'))
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            sp.remove(prstGeom)
        prstGeom = etree.SubElement(sp, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        corner_ratio = min(10 / min(aw, ah), 0.1) if min(aw, ah) > 0 else 0.02
        gd.set('fmla', f'val {int(corner_ratio * 50000)}')

        # Embed the video into the PPTX package
        _video_counter += 1
        ext = os.path.splitext(full)[1].lower()
        ct = VIDEO_CONTENT_TYPES.get(ext, 'video/mp4')

        with open(full, 'rb') as f:
            video_blob = f.read()

        from pptx.opc.packuri import PackURI
        try:
            from pptx.opc.part import Part
        except ImportError:
            from pptx.opc.package import Part

        partname = PackURI(f'/ppt/media/video{_video_counter}{ext}')
        video_part = Part(partname, ct, blob=video_blob, package=slide.part.package)

        RT_VIDEO = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/video'
        RT_MEDIA = 'http://schemas.microsoft.com/office/2007/relationships/media'

        vid_rId = slide.part.relate_to(video_part, RT_VIDEO)
        media_rId = slide.part.relate_to(video_part, RT_MEDIA)

        # Modify picture shape XML to add video references
        nvPicPr = pic._element.find(qn('p:nvPicPr'))
        cNvPr = nvPicPr.find(qn('p:cNvPr'))

        # Add hlinkClick for media playback action
        hlinkClick = etree.SubElement(cNvPr, qn('a:hlinkClick'))
        hlinkClick.set(qn('r:id'), '')
        hlinkClick.set('action', 'ppaction://media')

        # Add videoFile reference
        nvPr = nvPicPr.find(qn('p:nvPr'))
        videoFile = etree.SubElement(nvPr, qn('a:videoFile'))
        videoFile.set(qn('r:link'), vid_rId)

        # Add p14:media extension for PowerPoint 2010+
        extLst = etree.SubElement(nvPr, qn('p:extLst'))
        ext_el = etree.SubElement(extLst, qn('p:ext'))
        ext_el.set('uri', '{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}')

        p14_ns = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
        media_el = etree.SubElement(ext_el, f'{{{p14_ns}}}media')
        media_el.set(qn('r:embed'), media_rId)

        print(f"  Embedded video: {video_path}")

    except Exception as e:
        print(f"  Warning: video embed failed for {video_path}: {e}")
    finally:
        if poster_cleanup and os.path.exists(poster_cleanup):
            try:
                os.unlink(poster_cleanup)
            except OSError:
                pass


# ═══════════════════════════════════════════════════════════════
# Component builders (all coordinates in CSS pixels)
# ═══════════════════════════════════════════════════════════════

def _auto_scale_font(text, max_size, max_chars_at_max, min_size=18):
    """Reduce font size proportionally if text is too long for one line."""
    if len(text) <= max_chars_at_max:
        return max_size
    return max(min_size, int(max_size * max_chars_at_max / len(text)))


def add_heading(slide, text, y, is_lead=False, align=PP_ALIGN.LEFT):
    if is_lead:
        sz = _auto_scale_font(text, HEADING_LEAD_SIZE, HEADING_LEAD_CHARS, min_size=HEADING_LEAD_MIN)
        h = int(sz * 1.5)
        add_rich(slide, CL, y, CW, h, text, font=FONT_HEAD, size=sz, color=INK, bold=True, align=align)
        return y + h + 14
    # Content slides: push heading down
    y += 14
    sz = _auto_scale_font(text, HEADING_SIZE, HEADING_CHARS, min_size=HEADING_MIN)
    h = int(sz * 1.5)
    add_rich(slide, CL, y, CW, h, text, font=FONT_HEAD, size=sz, color=INK, bold=True)
    return y + h + 20




def _get_image_aspect(img_path):
    """Return width/height ratio of an image or video, or 1.0 if unreadable."""
    full = os.path.join(BASE_DIR, img_path)
    if _is_video(full):
        w, h = _get_video_dimensions(full)
        return w / h
    try:
        img = Image.open(full)
        return img.size[0] / img.size[1]
    except Exception:
        return 1.0


def add_image_grid(slide, images, y, grid_cls=2):
    if not images: return y
    n = len(images)
    gap = 8
    margin_top = 24
    cap_h = 24
    cap_gap = 6
    y += margin_top

    # Determine column count
    if n == 1:
        cols = 1
    elif grid_cls == 3:
        cols = 3
    elif grid_cls == 4:
        cols = 2
    else:
        cols = min(n, 2)

    total_rows = (n + cols - 1) // cols
    remaining = CANVAS_H - PAD_BOTTOM - y
    cell_w = (CW - gap * (cols - 1)) / cols

    # Get average aspect ratio to inform height calculation
    aspects = [_get_image_aspect(src) for src, _ in images]
    avg_aspect = sum(aspects) / len(aspects) if aspects else 1.0

    # Calculate ideal image height from cell width and aspect ratio
    ideal_h = cell_w / avg_aspect

    # Calculate max height that fits in remaining space
    max_h = (remaining - total_rows * (cap_gap + cap_h) - (total_rows - 1) * gap) / total_rows
    max_h = max(80, max_h)

    # Use the smaller of ideal and max, capped at 420px
    img_h = min(ideal_h, max_h, 420)
    img_h = max(80, img_h)
    row_h = img_h + cap_gap + cap_h

    for i, (src, caption) in enumerate(images):
        col = i % cols
        row = i // cols
        x = CL + col * (cell_w + gap)
        ry = y + row * (row_h + gap)

        if _is_video(src):
            add_video(slide, src, x, ry, cell_w, img_h)
        else:
            add_image(slide, src, x, ry, cell_w, img_h)
        add_text(slide, x, ry + img_h + cap_gap, cell_w, cap_h,
                 caption, size=CAPTION_FONT_SIZE, color=CAPTION_CLR, align=PP_ALIGN.CENTER)

    return y + total_rows * (row_h + gap)


def add_table(slide, headers, rows, y):
    if not headers or not rows: return y
    nc, nr = len(headers), len(rows) + 1
    margin = 8
    row_h = 36
    tbl_h = row_h * nr + margin * 2

    tbl_shape = slide.shapes.add_table(nr, nc, px(CL + margin), px(y + margin),
                                       px(CW - margin * 2), px(row_h * nr))
    tbl = tbl_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = re.sub(r'<[^>]+>', '', h)
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_BODY; p.font.size = Pt(TABLE_FONT_SIZE)
            p.font.color.rgb = TBL_HDR_TXT; p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = TBL_HDR_BG

    for i, row in enumerate(rows):
        cells = list(row) if isinstance(row, tuple) else [row]
        for j, val in enumerate(cells):
            if j < nc:
                cell = tbl.cell(i + 1, j)
                cell.text = re.sub(r'<[^>]+>', '', val)
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_BODY; p.font.size = Pt(TABLE_FONT_SIZE); p.font.color.rgb = TBL_BODY_TXT

    return y + tbl_h + 8


def add_chips(slide, chips, y, center=False, font_size=CHIP_FONT_SIZE):
    y += 8
    h = 26 if font_size >= 10 else 22
    gap = 8

    # Calculate total width of all chips for centering
    char_w = 8 if font_size >= 10 else 7
    widths = [max(70, len(text) * char_w + 24) for text in chips]
    total_w = sum(widths) + gap * (len(widths) - 1)
    x = CL + (CW - total_w) // 2 if center else CL

    for text, w in zip(chips, widths):
        if x + w > CANVAS_W - PAD_RIGHT:
            x = CL
            y += h + 4
        shp = add_card(slide, x, y, w, h, fill=CHIP_BG, border=CHIP_BORDER, shadow=False, radius=11)
        shp.text_frame.word_wrap = False
        shp.text_frame.margin_left = px(8)
        shp.text_frame.margin_right = px(8)
        shp.text_frame.margin_top = shp.text_frame.margin_bottom = px(2)
        p = shp.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY; p.font.size = Pt(font_size)
        p.font.color.rgb = CHIP_CLR; p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        x += w + gap

    return y + h + 8




def _estimate_bullet_lines(text, w):
    """Estimate how many visual lines a bullet wraps to in PPTX."""
    # Strip markdown formatting to get visible char count
    plain = re.sub(r'\*\*|`', '', text)
    # At BULLET_FONT_SIZE pt, avg char width ~0.8x font size in CSS px
    chars_per_line = max(1, int((w - 14) / (BULLET_FONT_SIZE * 0.8)))
    return max(1, -(-len(plain) // chars_per_line))  # ceil division


def add_bullets(slide, bullets, y, left=None, w=None):
    if not bullets: return y
    if left is None: left = CL
    if w is None: w = CW

    y += 8
    # Height per wrapped text line: 16pt * 1.15 line spacing = 18.4pt = 24.5 CSS px
    # Paragraph padding: space_before(3pt) + space_after(3pt) = 6pt = 8 CSS px
    wrapped_line_h = BULLET_LINE_H
    para_pad = BULLET_PARA_PAD
    total_h = 0
    for b in bullets:
        n_lines = _estimate_bullet_lines(b, w)
        total_h += n_lines * wrapped_line_h + para_pad

    txBox, tf = _make_txbox(slide, left, y, w, total_h)
    tf.margin_left = px(14)

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _fill_runs(p, f"\u2022  {bullet}", FONT_BODY, BULLET_FONT_SIZE, INK_SOFT, False)
        p.space_before = Pt(BULLET_SPACE_PT)
        p.space_after = Pt(BULLET_SPACE_PT)

    return y + total_h + 8


# ═══════════════════════════════════════════════════════════════
# Main slide builder
# ═══════════════════════════════════════════════════════════════

def build_slide(prs, text, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)

    is_lead = '<!-- _class: lead -->' in text
    heading = extract_heading(text)

    y = PAD_TOP

    # Lead slide: vertically + horizontally centered
    if is_lead:
        # Extract subtitle lines
        sub_lines = []
        after_h1 = False
        for line in text.split('\n'):
            s = line.strip()
            if s.startswith('# '): after_h1 = True; continue
            if after_h1 and s and not s.startswith('<') and not s.startswith('<!--'):
                sub_lines.append(s)
            if after_h1 and s.startswith('<div'): break

        chips = extract_chips(text)
        images = extract_images(text)

        # Calculate total content height for vertical centering
        content_h = 52 + 14  # heading + margin
        if sub_lines:
            content_h += 40 + 6  # subtitle + gap
        if chips:
            content_h += 8 + 22 + 8  # margin + chips + margin
        if images:
            content_h += 14 + 300 + 30  # margin + image + caption

        usable_h = CANVAS_H - PAD_TOP - PAD_BOTTOM
        y = PAD_TOP + max(0, (usable_h - content_h) // 2)

        if heading:
            y = add_heading(slide, heading, y, is_lead=True, align=PP_ALIGN.CENTER)
        if sub_lines:
            sub_text = ' '.join(sub_lines)
            sub_sz = _auto_scale_font(sub_text, SUBTITLE_SIZE, SUBTITLE_CHARS, min_size=SUBTITLE_MIN)
            add_rich(slide, CL, y, CW, 40, sub_text, size=sub_sz, color=INK_SOFT, align=PP_ALIGN.CENTER)
            y += 46
        if chips:
            y = add_chips(slide, chips, y, center=True, font_size=CHIP_LEAD_FONT_SIZE)
        if images:
            gc = extract_grid_class(text)
            y = add_image_grid(slide, images, y, grid_cls=gc)
        return

    if heading:
        y = add_heading(slide, heading, y, is_lead)

    bullets = extract_bullets(text)
    if bullets:
        y = add_bullets(slide, bullets, y)

    hdrs, rows = extract_table(text)
    if hdrs and rows:
        y = add_table(slide, hdrs, rows, y)

    images = extract_images(text)
    if images:
        gc = extract_grid_class(text)
        y = add_image_grid(slide, images, y, grid_cls=gc)


def parse_slides(md_path):
    with open(md_path, 'r') as f:
        content = f.read()
    content = re.sub(r'^---\n.*?\n---\n', '', content, count=1, flags=re.DOTALL)
    return [s.strip() for s in re.split(r'\n---\n', content)]


def main():
    parser = argparse.ArgumentParser(
        description='Convert Slide markdown to editable PPTX with CSS-matched styling')
    parser.add_argument('--input', '-i', required=True,
                        help='Input Slide markdown file')
    parser.add_argument('--output', '-o', required=True,
                        help='Output PPTX file path')
    parser.add_argument('--bg-image', '-b', required=True,
                        help='Background image file (e.g., slide_bg_light.png)')
    parser.add_argument('--base-dir', '-d', default='.',
                        help='Base directory for resolving image paths in markdown (default: cwd)')
    parser.add_argument('--theme', '-t', default='light', choices=['light', 'dark'],
                        help='Color theme for text/cards (light or dark)')
    parser.add_argument('--label', '-l', default=None,
                        help='Theme label for filename (e.g., warm). Defaults to --theme value.')
    args = parser.parse_args()

    apply_theme(args.theme)

    global BASE_DIR, BG_IMAGE_PATH
    BASE_DIR = os.path.abspath(args.base_dir)
    BG_IMAGE_PATH = os.path.abspath(args.bg_image)

    if not os.path.exists(args.input):
        print(f"Error: Input file not found: {args.input}")
        return 1
    if not os.path.exists(BG_IMAGE_PATH):
        print(f"Warning: Background image not found: {BG_IMAGE_PATH}")
        print("  Slides will have no background image.")

    # Pre-build: validate all image/video paths
    slides = parse_slides(args.input)
    missing = []
    for i, slide_text in enumerate(slides):
        for src, _ in extract_images(slide_text):
            full = os.path.join(BASE_DIR, src)
            if not os.path.exists(full):
                kind = 'video' if _is_video(src) else 'image'
                missing.append((i + 1, src, kind))
    if missing:
        print(f"\nERROR: {len(missing)} media file(s) not found:\n")
        for slide_num, src, kind in missing:
            print(f"  Slide {slide_num} ({kind}): {src}")
        print(f"\nBase dir: {BASE_DIR}")
        print("Fix the paths in the markdown and re-run.")
        return 1

    prs = Presentation()
    prs.slide_width = px(CANVAS_W)
    prs.slide_height = px(CANVAS_H)

    for i, slide_text in enumerate(slides):
        print(f"Building slide {i+1}...")
        build_slide(prs, slide_text, i)

    # Auto-generate output filename with date, time, and theme if output is a directory
    # or ends with slides.pptx (default name)
    output_path = args.output
    out_dir = os.path.dirname(output_path)
    out_base = os.path.basename(output_path)
    if out_base == 'slides.pptx' or out_base == '':
        timestamp = datetime.now().strftime('%m%d_%H%M')
        out_dir = out_dir or '.'
        label = args.label or args.theme
        output_path = os.path.join(out_dir, f'slides_{timestamp}_{label}.pptx')

    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    return 0


if __name__ == '__main__':
    exit(main() or 0)
